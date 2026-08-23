"""
File ingest — read any supported file, return (chunks, mime, page_count).

Dispatch is by extension. Each extractor returns a list[str] where each
element is one "page" of text; formats without pages (docx, xlsx, pptx,
txt) collapse to a single page or one page per slide/sheet.

OCR + vision extractors land in Phase 3c. This module handles the native-
text formats first so the rest of the pipeline (chunk → embed → store)
can be wired end to end today.
"""
from __future__ import annotations

import base64
from pathlib import Path
from typing import Awaitable, Callable

from chunker import Chunk, chunk_pages, chunk_text

# .pdf covers both native-text and scanned; scanned PDFs and raw images
# route through the VL model at ingest time (see needs_vision + chunks_for_vision).
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
SUPPORTED_EXTS = {".pdf", ".txt", ".md", ".docx", ".xlsx", ".pptx"} | IMAGE_EXTS

# Below this many characters across all extracted PDF pages we assume the PDF
# is scanned / image-only and fall back to VL OCR. Tuned to catch pages where
# pypdf returns a handful of garbage glyphs from a scan.
SCANNED_PDF_CHAR_THRESHOLD = 40

# Cap the number of pages we OCR per document to keep ingest bounded.
MAX_OCR_PAGES = 40


def _extract_pdf(path: Path) -> list[str]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages: list[str] = []
    for p in reader.pages:
        try:
            pages.append(p.extract_text() or "")
        except Exception:
            pages.append("")
    return pages


def _extract_txt(path: Path) -> list[str]:
    return [path.read_text(encoding="utf-8", errors="replace")]


def _extract_docx(path: Path) -> list[str]:
    from docx import Document
    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return ["\n".join(parts)]


def _extract_xlsx(path: Path) -> list[str]:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True, read_only=True)
    sheets: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = [f"# Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        sheets.append("\n".join(rows))
    return sheets  # one "page" per sheet


def _extract_pptx(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [f"# Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append("Notes: " + notes)
        slides.append("\n".join(parts))
    return slides  # one "page" per slide


EXTRACTORS: dict[str, Callable[[Path], list[str]]] = {
    ".pdf": _extract_pdf,
    ".txt": _extract_txt,
    ".md": _extract_txt,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def needs_vision(path: Path, pages: list[str] | None = None) -> bool:
    """
    True when the file must be routed through the VL model at ingest:
    - raw image formats
    - PDFs whose native text extraction is empty / scan-only
    """
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        return True
    if ext == ".pdf":
        if pages is None:
            try:
                pages = _extract_pdf(path)
            except Exception:
                return True
        total = sum(len((p or "").strip()) for p in pages)
        return total < SCANNED_PDF_CHAR_THRESHOLD
    return False


def extract(path: Path) -> tuple[list[str], str]:
    """Return (list_of_page_texts, mime_hint) for native-text formats only."""
    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"unsupported file type: {ext}")
    return EXTRACTORS[ext](path), ext.lstrip(".")


def chunks_for(path: Path) -> tuple[list[Chunk], str, int]:
    """Return (chunks, mime_hint, page_count) for a native-text file."""
    pages, mime = extract(path)
    if len(pages) <= 1:
        text = pages[0] if pages else ""
        return chunk_text(text, page=0), mime, 1 if text.strip() else 0
    return chunk_pages(pages), mime, len(pages)


async def chunks_via_vision(
    path: Path,
    extract_text: Callable[[str], Awaitable[str]],
) -> tuple[list[Chunk], str, int]:
    """
    Route an image or scanned PDF through a VL model. `extract_text` is an
    async callable (b64_png) -> extracted_text so this module does not need
    to know about ollama_client directly (easier to test, easier to swap for
    a local OCR engine later).
    """
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        b64 = base64.b64encode(path.read_bytes()).decode("ascii")
        text = await extract_text(b64)
        return chunk_text(text, page=1), ext.lstrip("."), 1

    if ext == ".pdf":
        import fitz  # PyMuPDF
        pages_text: list[str] = []
        doc = fitz.open(str(path))
        try:
            for i, page in enumerate(doc):
                if i >= MAX_OCR_PAGES:
                    break
                pix = page.get_pixmap(dpi=150)
                png = pix.tobytes("png")
                b64 = base64.b64encode(png).decode("ascii")
                text = await extract_text(b64)
                pages_text.append(text)
        finally:
            doc.close()
        return chunk_pages(pages_text), "pdf", len(pages_text)

    raise ValueError(f"vision path does not support {ext}")
